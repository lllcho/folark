"""PPTX 文件处理实现。"""

from __future__ import annotations

import logging
from typing import TYPE_CHECKING

from app.plugins.builtin_plugin.utils import render_template
from app.plugins.core import PreviewResult

if TYPE_CHECKING:
    from app.plugins.core import PluginContext

logger = logging.getLogger(__name__)


def extract(ctx: PluginContext) -> str | None:
    """从 PPTX 文件提取文本内容。"""
    try:
        from pptx import Presentation

        prs = Presentation(str(ctx.file_path))
        texts: list[str] = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            texts.append(text)
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            text = cell.text.strip()
                            if text:
                                texts.append(text)
        return "\n".join(texts) if texts else None
    except Exception as e:
        logger.warning("PPTX 文本提取失败 %s: %s", ctx.file_path, e)
        return None


def convert(ctx: PluginContext) -> bytes | None:
    """将 PPTX 文件转换为 PDF。"""
    if ctx.target_type != "pdf":
        logger.warning("PPTX 不支持转换为 %s 格式", ctx.target_type)
        return None

    try:
        import pymupdf

        doc = pymupdf.open(ctx.file_path)
        pdf_bytes = doc.convert_to_pdf()
        doc.close()

        return pdf_bytes
    except Exception as e:
        logger.warning("PPTX 转 PDF 失败 %s: %s", ctx.file_path, e)
        return None


def preview(ctx: PluginContext) -> PreviewResult:
    """使用模板渲染 PPTX 预览页面。"""
    html = render_template(
        "pptx_preview.html",
        {
            "file_name": ctx.file_name,
            "title": ctx.title or ctx.file_name,
            "file_url": ctx.file_url,
        },
    )
    return PreviewResult.from_html(html)
